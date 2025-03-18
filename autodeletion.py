import os
import sys
import time
import multiprocessing
from PyPDF2 import PdfReader
from PIL import Image, ImageTk
import imagehash
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import textract
from docx import Document
import re
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import PorterStemmer
from nltk.stem import WordNetLemmatizer
from pptx import Presentation
import csv
import json
import openpyxl
import xml.etree.ElementTree as ET
import markdown2 
# from kivy.app import App
# from kivy.uix.boxlayout import BoxLayout
# from kivy.uix.button import Button
# from kivy.uix.label import Label
# from kivy.uix.textinput import TextInput
# Initialize NLTK resources
stop_words = set(stopwords.words('english'))
ps = PorterStemmer()
lemmatizer = WordNetLemmatizer()

# Dictionary to store preprocessed text content of each file
preprocessed_text_cache = {}
class CheckboxTreeview(tk.Frame):
    def __init__(self, master, **kwargs):
        super().__init__(master, **kwargs)
        self.checked_items = set()
        self.treeview = ttk.Treeview(self, selectmode="extended", columns=("Filename", "Similarity"))
        self.treeview.heading("#0", text="Check")
        self.treeview.heading("Filename", text="Filename")
        self.treeview.heading("Similarity", text="Similarity")
        self.treeview.column("#0", width=50)  # Adjust the width of the Check column if needed
        self.treeview.column("Filename", width=500)  # Adjust the width of the Filename column as needed
        self.treeview.column("Similarity", width=100)  # Adjust the width of the Similarity column as needed
        self.treeview.place(x=200,y=200)
        self.treeview.pack(side=tk.BOTTOM, fill=tk.BOTH, expand=True)
        checked_img = Image.open("checked.png")
        checked_img = checked_img.resize((16, 16))  # Resize to 16x16 pixels

        # Convert resized images to PhotoImage
        self.checked_image = ImageTk.PhotoImage(checked_img)

        self.treeview.bind("<Button-1>", self.toggle_checkbox)

    def insert_checkbox(self, index, filename, similarity):
        item_id = self.treeview.insert("", index, text="", values=(filename, similarity))
        self.checked_items.discard(item_id)

    def toggle_checkbox(self, event):
        item_id = self.treeview.identify_row(event.y)
        if item_id:
            if item_id in self.checked_items:
                self.treeview.item(item_id, image="")
                self.checked_items.remove(item_id)
            else:
                self.treeview.item(item_id, image=self.checked_image)
                self.checked_items.add(item_id)

    def get_checked_items(self):
        checked_items = []
        for item_id in self.checked_items:
            try:
                item_values = self.treeview.item(item_id)
                if item_values is not None:  # Check if the item exists in the treeview
                    checked_items.append(item_values["values"][0])
            except tk.TclError:  # Catch the TclError if the item is not found
                pass
        return checked_items
    def select_items_above_similarity(self, threshold):
        if self.checked_items:  # If some items are already checked, deselect all
            self.deselect_all_items()
          # If no items are checked, select all above the threshold
        for item_id in self.treeview.get_children():
            similarity = float(self.treeview.item(item_id, "values")[1][:-1])
            if similarity >= threshold:
                self.treeview.item(item_id, image=self.checked_image)
                self.checked_items.add(item_id)
    def deselect_all_items(self):
        for item_id in self.checked_items.copy():  # Use copy() to avoid modifying set while iterating
            try:
                self.treeview.item(item_id, image="")
            except tk.TclError:
                print(f"Item {item_id} not found.")
            else:
                self.checked_items.remove(item_id)  # Remove the item ID if it was successfully deselected



def rolling_hash(text, window_size):
    """
    Compute rolling hash values for all windows of size `window_size`.
    """
    hash_values = []
    text_len = len(text)
    prime = 101  # Choose a prime number
    modulus = 2**32  # Typically a large prime number
    hash_value = 0
    for i in range(window_size):
        hash_value = (hash_value * prime + ord(text[i])) % modulus
    hash_values.append(hash_value)

    for i in range(1, text_len - window_size + 1):
        hash_value = (hash_value * prime - ord(text[i - 1]) * pow(prime, window_size, modulus) + ord(text[i + window_size - 1])) % modulus
        hash_values.append(hash_value)
    return hash_values

def preprocess_text(text):
    # Remove special characters, punctuation, and extra whitespaces
    text = re.sub(r'[^A-Za-z0-9\s]', '', text)
    text = re.sub(r'\s+', ' ', text)
    # Convert text to lowercase
    text = text.lower()
    # Tokenize the text
    words = word_tokenize(text)
    # Remove stopwords
    words = [word for word in words if word not in stop_words]
    # Lemmatize words (or alternatively, use stemming with ps.stem(word))
    words = [lemmatizer.lemmatize(word) for word in words] #root form ing,ed remove
    # Join the preprocessed words back into a single string
    preprocessed_text = ' '.join(words)
    return preprocessed_text
def extract_text_from_csv(csv_path):
    try:
        with open(csv_path, 'r', newline='', encoding='utf-8') as csvfile:
            reader = csv.reader(csvfile)
            text = '\n'.join(','.join(row) for row in reader)
        
        # Preprocess the extracted text
        preprocessed_text = preprocess_text(text)
        return preprocessed_text
    except Exception as e:
        print(f"Error extracting text from {csv_path}: {e}")
        return None

# Function to extract text from JSON files
def extract_text_from_json(json_path):
    try:
        with open(json_path, 'r', encoding='utf-8') as jsonfile:
            data = json.load(jsonfile)
            # Convert JSON data to text (customize as per your JSON structure)
            text = '\n'.join([str(item) for item in data])
        
        # Preprocess the extracted text
        preprocessed_text = preprocess_text(text)
        return preprocessed_text
    except Exception as e:
        print(f"Error extracting text from {json_path}: {e}")
        return None

# Function to extract text from XLSX files
def extract_text_from_xlsx(xlsx_path):
    try:
        workbook = openpyxl.load_workbook(xlsx_path)
        text = ''
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                text += ' '.join(str(cell) for cell in row) + '\n'
        
        # Preprocess the extracted text
        preprocessed_text = preprocess_text(text)
        return preprocessed_text
    except Exception as e:
        print(f"Error extracting text from {xlsx_path}: {e}")
        return None

# Function to extract text from XML files
def extract_text_from_xml(xml_path):
    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
        # Extract text from XML structure (customize as per your XML structure)
        text = ' '.join(element.text for element in root.iter() if element.text)
        
        # Preprocess the extracted text
        preprocessed_text = preprocess_text(text)
        return preprocessed_text
    except Exception as e:
        print(f"Error extracting text from {xml_path}: {e}")
        return None

# Function to extract text from Markdown (.md) files
def extract_text_from_markdown(md_path):
    try:
        with open(md_path, 'r', encoding='utf-8') as mdfile:
            text = mdfile.read()
        # Convert Markdown to plain text
        text = markdown2.markdown(text)
        
        # Preprocess the extracted text
        preprocessed_text = preprocess_text(text)
        return preprocessed_text
    except Exception as e:
        print(f"Error extracting text from {md_path}: {e}")
        return None
def extract_text_from_ppt(ppt_path):
    try:
        presentation = Presentation(ppt_path)
        text = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        
        # Preprocess the extracted text
        preprocessed_text = preprocess_text(text)
        return preprocessed_text
    except Exception as e:
        print(f"Error extracting text from {ppt_path}: {e}")
        return None

def extract_text_from_pdf(pdf_path):
    print(f"Reading PDF: {pdf_path}")
    with open(pdf_path, 'rb') as f:
        reader = PdfReader(f)
        text = ''
        total_pages = len(reader.pages)
        for i, page in enumerate(reader.pages, 1):
            print(f"Processing page {i}/{total_pages}")
            text += page.extract_text()
        print("PDF reading complete.")     
    # Preprocess the extracted text
    preprocessed_text = preprocess_text(text)
    return preprocessed_text
def extract_text_from_docx(docx_path):
    try:
        doc = Document(docx_path)
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text
        
        # Preprocess the extracted text
        preprocessed_text = preprocess_text(text)
        return preprocessed_text
    except Exception as e:
        print(f"Error extracting text from {docx_path}: {e}")
        return None
def extract_text_from_file(file_path):
    """
    Extract text content from a supported file format.
    """
    supported_extensions = ['.ppt', '.pptx','.txt', '.pdf', '.doc', '.docx', '.rtf', '.html', '.htm', '.odt','.csv','.json','.xlsx','.xml','.md']  # Add more extensions as needed
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == '.pdf':
        # Extract text from PDF
        return extract_text_from_pdf(file_path)
    elif file_extension == '.txt':
        # Read text from plain text file
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    elif file_extension in ['.doc','.docx']:
        text2 = extract_text_from_docx(file_path)
        return text2
    elif file_extension in ['.ppt', '.pptx']:
    # Extract text from PowerPoint presentations
        return extract_text_from_ppt(file_path)
    elif file_extension == '.csv':
        return extract_text_from_csv(file_path)
    elif file_extension == '.json':
        return extract_text_from_json(file_path)
    elif file_extension == '.xlsx':
        return extract_text_from_xlsx(file_path)
    elif file_extension == '.xml':
        return extract_text_from_xml(file_path)
    elif file_extension == '.md':
        return extract_text_from_markdown(file_path)
    elif file_extension in ['.rtf', '.html', '.htm', '.odt']:
        # Extract text from other supported formats using textract
        try:
            text = textract.process(file_path, encoding='utf-8').decode('utf-8')
            return text
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return ""
    else:
        print(f"Unsupported file format: {file_extension}")
        return ""


def find_similarity(file1, file2, window_size):
    """
    Find similarity between two text files using rolling hashing.
    """
    # print(f"Calculating similarity between {file1} and {file2}")
    text1 = preprocessed_text_cache.get(file1)
    if text1 is None:
        text1 = extract_text_from_file(file1)
        preprocessed_text_cache[file1] = text1

    text2 = preprocessed_text_cache.get(file2)
    if text2 is None:
        text2 = extract_text_from_file(file2)
        preprocessed_text_cache[file2] = text2

    hash_values1 = set(rolling_hash(text1, window_size))
    hash_values2 = set(rolling_hash(text2, window_size))

    common_hashes = hash_values1.intersection(hash_values2)
    similarity = (len(common_hashes) / (len(hash_values1) + len(hash_values2) - len(common_hashes))) * 100

    # print(f"Similarity calculation complete.")
    return similarity
def calculate_image_hash(image_path):
    """
    Calculate the hash value of an image.
    """
    try:
        with Image.open(image_path) as img:
            # Convert the image to grayscale and calculate its hash
            hash_value = imagehash.phash(img)
        return str(hash_value)
    except Exception as e:
        print(f"Error calculating hash for {image_path}: {e}")
        return None

def find_similar_images(image_path, directory):
    """
    Find similar images in the directory compared to the given image.
    """
    image_hash = calculate_image_hash(image_path)
    if image_hash is None:
        print("Error: Unable to calculate hash for the image.")
        return []

    similar_images = []
    for root, _, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            if file_path.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff')):
                hash_value = calculate_image_hash(file_path)
                if hash_value is not None:
                    # Calculate similarity using Hamming distance
                    similarity = 1 - (imagehash.hex_to_hash(image_hash) - imagehash.hex_to_hash(hash_value)) / 64
                    similar_images.append((file_path, similarity * 100))
    return similar_images
def find_related_files(directory, file1, window_size):
    """
    Find and compare similarity of related files in the directory and its subdirectories to file1.
    """
    related_files = []
    new_file_size = os.path.getsize(file1)
    supported_extensions = ['.txt', '.pdf', '.doc', '.docx', '.rtf', '.html', '.htm', '.odt','.csv','.json','.xlsx','.xml','.md','.ppt', '.pptx']  # Add more extensions as needed

    # Iterate over all directories, subdirectories, and files in the directory tree
    for root, dirs, files in os.walk(directory):
        for filename in files:
            filepath = os.path.join(root, filename)
            if os.path.splitext(filename)[1].lower() in supported_extensions: #if filepath != file1 and os.path.splitext(filename)[1].lower() in supported_extensions:
                if os.path.getsize(filepath) <= new_file_size + new_file_size/2:
                    try:
                        similarity = find_similarity(file1, filepath, window_size)
                        related_files.append((filepath, similarity)) #filename to filepath
                    except Exception as e:
                        print(f"Error processing {filename}: {e}")

    return related_files


def callSimilar(file_path, directory, window_size):
    file1 = file_path
    related_files = find_related_files(directory, file1, window_size)
    print(f"Similarity of {file1} with other files in the {directory}:")
    for filename, similarity in related_files:
        # print(f"{filename}: {similarity:.2f}%")
        pass
    return related_files

class NewFileHandler(FileSystemEventHandler):
    def __init__(self, directory, window_size, delay=5):
        self.directory = directory
        self.window_size = window_size
        self.delay = delay  # Delay in seconds
    def on_created(self, event):
        if event.is_directory:
            return
        file_path = event.src_path
        if file_path.endswith('.crdownload') or not file_path.lower().endswith(('.txt', '.pdf', '.doc', '.docx', '.rtf', '.html', '.htm', '.odt','.csv','.json','.xlsx','.xml','.md','.ppt', '.pptx','.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff')):
            print(f"Skipping non-relevant file: {file_path}")
            return
        time.sleep(self.delay)
        var=None
        if file_path.endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff')):
            output =find_similar_images(file_path, self.directory)
            var=1
            print(output)
        else:
            output = callSimilar(file_path, self.directory, self.window_size)
            var=2
        temp_var=0
        if var==1:
            for _, simil in output:
                if simil > 85 and _ != file_path:
                    temp_var=1
                    break
        if var==2:
            for _, simil in output:
                if simil > 10 and _ != file_path:
                    temp_var=1
                    break
        if temp_var:
            popup_window(output, file_path,var)
        else:
            pass #pop message if new file with no similairty or empty file
        print(f"New file created: {file_path}")
        # popup_window(output, file_path)

    # def on_modified(self, event):
    #     if event.is_directory:
    #         return
    #     file_path = event.src_path
    #     time.sleep(self.delay)  # Introduce a delay before processing the file
    #     output = callSimilar(file_path, self.directory, self.window_size)
    #     print(f"File modified: {file_path}")
    #     popup_window(output, file_path)

def popup_window(files, file_path,var, autocheck=True):
    if not files or all(similarity <= 10 for _, similarity in files):  # If no similar files found or all are <= 10%
        print("No similar files found or similarity <= 10%.")
        return
    def delete_selected_files():
        nonlocal files
        checked_items = checkbox_treeview.get_checked_items()
        files_to_delete = []  # Accumulate files to delete
        for item in checked_items:
            selected_filename = item
            selected_file = next((file for file in files if file[0] == selected_filename), None)
            if selected_file:
                files_to_delete.append(selected_file)  # Accumulate the file object
        if files_to_delete:
            # Construct a confirmation message listing all files to be deleted
            confirmation_message = f"Are you sure you want to delete the following files?\n\n"
            confirmation_message += "\n".join(file[0] for file in files_to_delete)
            # Display a confirmation messagebox
            if messagebox.askyesno("Confirm Deletion", confirmation_message):
                # Delete files and display a success message
                for file_to_delete in files_to_delete:
                    files.remove(file_to_delete)
                    filepath = file_to_delete[0]
                    try:
                        os.remove(filepath)
                    except OSError as e:
                        messagebox.showerror("Error", f"Failed to delete {file_to_delete[0]}: {e}")
                messagebox.showinfo("Files Deleted", "Selected files have been deleted successfully.")
                refresh_treeview()
        elif checked_items:  # Add this condition to check if any files are selected for deletion
            messagebox.showinfo("No Files Selected", "No files selected for deletion.")

    def auto_delete_selected_files():
        nonlocal files
        checked_items = checkbox_treeview.get_checked_items()
        files_to_delete = []
        for item in checked_items:
            selected_filename = item
            selected_file = next((file for file in files if file[0] == selected_filename), None)
            if selected_file:
                files.remove(selected_file)
                files_to_delete.append(selected_filename)
        auto_checked_for_deletion = [filename for filename in autochecked_files if filename not in files_to_delete]
        if auto_checked_for_deletion:
                # Delete auto-checked files and display a success message
            for filename in auto_checked_for_deletion:
                filepath =filename
                try:
                    os.remove(filepath)
                except OSError as e:
                    messagebox.showerror("Error", f"Failed to delete {filename}: {e}")
            messagebox.showinfo("Files Deleted", "Auto-checked files with 100% similarity have been deleted successfully.")
            # Remove auto-checked files from the list before refreshing the treeview
            files = [file for file in files if file[0] not in auto_checked_for_deletion]
            refresh_treeview()

    def refresh_treeview():
        checkbox_treeview.treeview.delete(*checkbox_treeview.treeview.get_children())
        for filename, similarity in files:
            cur_file = os.path.basename(file_path)
            if filename != cur_file:
                checkbox_treeview.insert_checkbox(tk.END, filename, f"{similarity:.2f}%")
    # def toggle_select90():
    #     if select_90_var.get():
    #         checkbox_treeview.select_items_above_similarity(90)
    #     elif not select_80_var.get():
    #         checkbox_treeview.deselect_all_items()
    # def toggle_select80():
    #     if select_80_var.get():
    #         checkbox_treeview.select_items_above_similarity(80)
    #     elif not select_90_var.get():
    #         checkbox_treeview.deselect_all_items()
    root = tk.Tk()
    root.title("File Similarity Checker")
    root.geometry("800x500")
    root.attributes('-topmost', True)
    main_frame = tk.Frame(root)
    main_frame.pack(pady=5)
    checkbox_treeview = CheckboxTreeview(main_frame)
    checkbox_treeview.place(x=200,y=200)
    checkbox_treeview.pack(side=tk.BOTTOM, fill=tk.BOTH, expand=True)
    item_ids = {}  # Dictionary to store item IDs corresponding to filenames
    autochecked_files = []  
    matching = ""  # Files to be auto-checked
    checkIf = 0
    new_file=file_path
    oldest_files = find_old_files(files)
    first_file, first_similarity = next(iter(oldest_files.items()))
    # print(first_file,first_similarity,"found")
    label = tk.Label(root, text=f"New file detected {new_file}.")  # Update label text to display file_path
    label.pack(pady=5)
    if (first_file and first_similarity is not None) and first_similarity==100:
        file_path=first_file
        label = tk.Label(root, text=f" The following files having similarity with {os.path.basename(file_path)}")  # Update label text to display file_path
        label.pack()
    for filename, similarity in files:
        cur_file = os.path.basename(file_path)
        # print(filename,cur_file) # Assuming file_path is defined elsewhere
        if os.path.basename(filename) != cur_file and similarity>10 and var==2: 
             # Correct usage of os.path.dirname()
            item_id = checkbox_treeview.insert_checkbox(tk.END, filename, f"{similarity:.2f}%")
            item_ids[filename] = item_id  # Store item ID corresponding to filename
            if autocheck:  # Assuming autocheck is defined elsewhere
                if similarity == 100 and (checkIf == 1 or checkIf == 0):
                    autochecked_files.append(filename)
                    matching = "100"
                    checkIf = 1
                elif similarity >= 90 and (checkIf == 2 or checkIf == 0):
                    autochecked_files.append(filename)
                    matching = "Above 90"
                    checkIf = 2
                elif similarity >= 80 and (checkIf == 3 or checkIf == 0):
                    autochecked_files.append(filename)
                    matching = "Above 80"
                    checkIf = 3
        if os.path.basename(filename) != cur_file and similarity>80 and var==1: 
             # Correct usage of os.path.dirname()
            item_id = checkbox_treeview.insert_checkbox(tk.END, filename, f"{similarity:.2f}%")
            item_ids[filename] = item_id  # Store item ID corresponding to filename
            if autocheck:  # Assuming autocheck is defined elsewhere
                if similarity == 100 and (checkIf == 1 or checkIf == 0):
                    autochecked_files.append(filename)
                    matching = "100"
                    checkIf = 1
                elif similarity >= 90 and (checkIf == 2 or checkIf == 0):
                    autochecked_files.append(filename)
                    matching = "Above 90"
                    checkIf = 2
                elif similarity >= 80 and (checkIf == 3 or checkIf == 0):
                    autochecked_files.append(filename)
                    matching = "Above 80"
                    checkIf = 3

    delete_button = tk.Button(root, text="Delete Selected Files", command=delete_selected_files)
    delete_button.pack(pady=5)

    def on_slider_move(event):
        similarity_threshold = similarity_slider.get()   # Multiply by 2 to get even numbers
        checkbox_treeview.select_items_above_similarity(similarity_threshold)
    style = ttk.Style()

    similarity_slider = tk.Scale(root, from_=0, to=100, orient=tk.HORIZONTAL, length=200, command=on_slider_move)
    similarity_slider.set(0)  
    similarity_slider.pack(pady=5)
    # Your delete button and slider code
    if autocheck and autochecked_files:
        # Automatically select files with high similarity
        for filename in autochecked_files:
            item_id = item_ids.get(filename)
            if item_id:
                checkbox_treeview.treeview.selection_set(item_id)

        confirmation_message = f"New file detected {new_file}. The following files have {matching}% similarity with {os.path.basename(file_path)}. Do you want to delete them?\n\n"
        confirmation_message += "\n".join(autochecked_files)
        if messagebox.askyesno("Confirm Deletion", confirmation_message):
            auto_delete_selected_files()  # Delete files if confirmed
    root.mainloop()
def find_old_files(file_info):
    oldest_files = {}
    oldest_creation_time = float('inf')  # Initialize with infinity
    
    # Iterate through each file path and similarity score in the dictionary
    for file_path, sim_score in file_info:
        # Get the creation time of the file
        creation_time = os.path.getctime(file_path)
        
        # Check if the current file is older than the previously found oldest file(s)
        if creation_time < oldest_creation_time and sim_score==100:
            # Clear the previous dictionary of oldest files
            oldest_files = {file_path: sim_score}
            # Update the oldest creation time
            oldest_creation_time = creation_time
        elif creation_time == oldest_creation_time and sim_score==100:
            # If there are multiple files with the same oldest creation time, add them to the dictionary
            oldest_files[file_path] = sim_score
    
    return oldest_files

# popup_window([("file1.txt", 90), ("file2.txt", 80)], "path/to/current/file.txt")

def watch_directory(directory, window_size):
    event_handler = NewFileHandler(directory, window_size)
    observer = Observer()
    observer.schedule(event_handler, directory, recursive=True)
    observer.start()
    # print("c6")
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()

def start_background_process(directory, window_size):
    process = multiprocessing.Process(target=watch_directory, args=(directory, window_size))
    process.daemon = True  # Make the process a daemon so it runs in the background
    process.start()
    return process

def stop_background_process(process):
    process.terminate()

def start():
    global background_process
    # directory = 'D:/'
    # window_size = 8
    directory = directory_entry.get()
    window_size = int(window_size_entry.get())
    if not os.path.isdir(directory):
        messagebox.showerror("Error", f"{directory} is not a valid directory.")
        return
    # if background_process and background_process.is_alive():
    #     messagebox.showwarning("Warning", "A process is already running.")
    #     return
    background_process = start_background_process(directory, window_size)
    print("Background process started.")
    start_button.config(state=tk.DISABLED)
    # print("c5")

def stop():
    global background_process
    if background_process:
        stop_background_process(background_process)
        print("Background process stopped.")
    else:
        print("No background process to stop.")
    start_button.config(state=tk.NORMAL)
# added_files = [
#     ('checked.png', '.'),  # Add the image file
# ]
def browse_directory():
    selected_directory = filedialog.askdirectory()
    directory_entry.delete(0, tk.END)
    directory_entry.insert(0, selected_directory)
if __name__ == "__main__":
    root = tk.Tk()
    root.title("File Similarity Checker")
    root.geometry("400x200")

    directory_label = tk.Label(root, text="Directory:")
    directory_label.pack()

    directory_entry = tk.Entry(root)
    directory_entry.pack()

    browse_button = tk.Button(root, text="Browse", command=browse_directory)
    browse_button.pack()

    window_size_label = tk.Label(root, text="Window Size:")
    window_size_label.pack()

    window_size_entry = tk.Entry(root)
    window_size_entry.pack()

    start_button = tk.Button(root, text="Start", command=start)
    start_button.pack()

    stop_button = tk.Button(root, text="Stop", command=stop)
    stop_button.pack()

    background_process = None
    # print("c1")
    root.mainloop()
# if __name__ == "__main__":
#     if len(sys.argv) != 3:
#         print("Usage: python script.py /path/to/directory window_size")
#         sys.exit(1)

#     directory = sys.argv[1]
#     window_size = int(sys.argv[2])
#     print(directory)
#     if not os.path.isdir(directory):
#         print(f"Error: {directory} is not a directory.")
#         sys.exit(1)
#     watch_directory(directory, window_size)

# from kivy.app import App
# from kivy.uix.boxlayout import BoxLayout
# from kivy.uix.button import Button
# from kivy.uix.label import Label
# from kivy.uix.textinput import TextInput
# class MainWindow(BoxLayout):
#     def __init__(self, **kwargs):
#         super().__init__(**kwargs)
#         self.orientation = "vertical"
#         self.directory_label = Label(text="Directory:")
#         self.directory_input = TextInput()
#         self.window_size_label = Label(text="Window Size:")
#         self.window_size_input = TextInput()
#         self.run_button = Button(text="Run", on_press=self.run_similarity_check)
#         self.pause_button = Button(text="Pause", on_press=self.pause_similarity_check)
#         self.add_widget(self.directory_label)
#         self.add_widget(self.directory_input)
#         self.add_widget(self.window_size_label)
#         self.add_widget(self.window_size_input)
#         self.add_widget(self.run_button)
#         self.add_widget(self.pause_button)
#         self.running = True
#         self.event_handler = None

#     def run_similarity_check(self, instance):
#         if not self.running:
#             self.resume_similarity_check()
#             return
#         directory = self.directory_input.text
#         window_size = int(self.window_size_input.text)
#         if not os.path.isdir(directory):
#             print(f"Error: {directory} is not a directory.")
#             return
#         self.running = True
#         self.event_handler = NewFileHandler(directory, window_size)
#         observer = Observer()
#         observer.schedule(self.event_handler, directory, recursive=True)
#         observer.start()

#     def pause_similarity_check(self, instance):
#         self.running = False
#         if self.event_handler:
#             self.event_handler.pause()

#     def resume_similarity_check(self):
#         self.running = True
#         if self.event_handler:
#             self.event_handler.resume()

# class FileSimilarityCheckerApp(App):
#     def build(self):
#         return MainWindow()


# if __name__ == "__main__":
#     FileSimilarityCheckerApp().run()